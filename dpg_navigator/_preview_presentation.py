"""PowerPoint document loading for the preview panel.

Extracts slide content without depending on DearPyGui or Pillow.
"""

from __future__ import annotations

# MIT licensed
from dataclasses import dataclass
from typing import Any, Callable

from ._preview_registry import ooxml_exceeds_preview_limit

_Presentation: Any
try:
    from pptx import Presentation as _Presentation  # type: ignore[import-untyped]
except Exception:  # optional backend absent or incompatible (e.g. old Python)
    _Presentation = None


class PresentationPreviewError(Exception):
    """PowerPoint presentation data could not be loaded."""


@dataclass(frozen=True)
class PresentationRun:
    """Inline text with the formatting used by the preview panel."""

    text: str
    bold: bool
    italic: bool


@dataclass(frozen=True)
class PresentationParagraph:
    """Text paragraph extracted from a slide shape."""

    text: str
    level: int
    runs: list[PresentationRun]


@dataclass(frozen=True)
class PresentationTable:
    """Table rows extracted from a slide shape."""

    rows: list[list[str]]


@dataclass(frozen=True)
class PresentationShape:
    """Preview-ready shape content."""

    table: PresentationTable | None
    image_blob: bytes | None
    paragraphs: list[PresentationParagraph]


@dataclass(frozen=True)
class PresentationSlide:
    """Preview-ready slide content."""

    shapes: list[PresentationShape]
    notes: str


@dataclass(frozen=True)
class PresentationDocument:
    """Preview-ready presentation content."""

    slides: list[PresentationSlide]


def presentation_available() -> bool:
    """Return True when PowerPoint document loading is available."""
    return _Presentation is not None


def _load_notes(slide: Any) -> str:
    try:
        notes_slide = slide.notes_slide
        if notes_slide and notes_slide.notes_text_frame:
            return str(notes_slide.notes_text_frame.text.strip())
    except Exception:
        pass
    return ""


def _load_shape(shape: Any) -> PresentationShape:
    if shape.has_table:
        return PresentationShape(
            table=PresentationTable([[cell.text.strip() for cell in row.cells] for row in shape.table.rows]),
            image_blob=None,
            paragraphs=[],
        )

    image_blob = None
    if shape.shape_type is not None and hasattr(shape, "image"):
        try:
            image_blob = shape.image.blob
        except Exception:
            pass

    paragraphs = []
    if shape.has_text_frame:
        for paragraph in shape.text_frame.paragraphs:
            paragraphs.append(
                PresentationParagraph(
                    text=paragraph.text,
                    level=paragraph.level or 0,
                    runs=[
                        PresentationRun(
                            run.text,
                            bool(run.font.bold),
                            bool(run.font.italic),
                        )
                        for run in paragraph.runs
                    ],
                )
            )

    return PresentationShape(
        table=None,
        image_blob=image_blob,
        paragraphs=paragraphs,
    )


def load_presentation(
    path: str,
    *,
    presentation_loader: Callable[..., Any] | None = _Presentation,
) -> PresentationDocument:
    """Load slides, shapes, and speaker notes in display order."""
    if presentation_loader is None:
        raise PresentationPreviewError("python-pptx is not installed")
    if ooxml_exceeds_preview_limit(path):
        raise PresentationPreviewError("File too large for preview")

    try:
        presentation = presentation_loader(path)
        slides = [
            PresentationSlide(
                shapes=[_load_shape(shape) for shape in slide.shapes],
                notes=_load_notes(slide),
            )
            for slide in presentation.slides
        ]
    except Exception as exc:
        raise PresentationPreviewError(str(exc)) from exc

    return PresentationDocument(slides)
